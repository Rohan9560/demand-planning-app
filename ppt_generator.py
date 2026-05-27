from pptx import Presentation
from pptx.util import Inches
import uuid
import os

def create_ppt(data):

    os.makedirs("generated_reports", exist_ok=True)

    prs = Presentation()

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    slide.shapes.title.text = "Demand Planning Assessment"

    textbox = slide.shapes.add_textbox(
        Inches(1),
        Inches(1.5),
        Inches(6),
        Inches(2)
    )

    textbox.text_frame.text = (
        f"Company: {data.get('company', 'N/A')}\n"
        f"Score: {data.get('score', 'N/A')}\n"
        f"Stage: {data.get('stage', 'N/A')}"
    )

    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])

    slide2.shapes.title.text = "Assessment Summary"

    textbox2 = slide2.shapes.add_textbox(
        Inches(1),
        Inches(1.5),
        Inches(6),
        Inches(3)
    )

    textbox2.text_frame.text = (
        f"People Score: {data.get('people', 'N/A')}\n"
        f"Process Score: {data.get('process', 'N/A')}\n"
        f"Data Score: {data.get('data', 'N/A')}\n"
        f"Technology Score: {data.get('technology', 'N/A')}"
    )

    filename = f"generated_reports/{uuid.uuid4()}.pptx"

    prs.save(filename)

    print("PPT Saved:", filename)

    return filename